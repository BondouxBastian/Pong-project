#!/usr/bin/env python3
"""Generates pong-deck.pptx from the content of pong-deck.html (dark VT323/IBM Plex look)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BG = RGBColor(0x0B, 0x0C, 0x0A)
SURFACE = RGBColor(0x16, 0x17, 0x0F)
SURFACE2 = RGBColor(0x1E, 0x20, 0x15)
LINE = RGBColor(0x35, 0x34, 0x2A)
LINE_STRONG = RGBColor(0x55, 0x52, 0x3F)
TEXT = RGBColor(0xF4, 0xEF, 0xE2)
TEXT_DIM = RGBColor(0xA3, 0x9D, 0x8A)
TEXT_FAINT = RGBColor(0x6E, 0x6A, 0x5B)
ACCENT = RGBColor(0xFF, 0x8F, 0x2E)
ACCENT_SOFT = RGBColor(0x7A, 0x4A, 0x1E)
INFO = RGBColor(0x6F, 0xB1, 0xE0)

FONT_BODY = "IBM Plex Sans"
FONT_MONO = "IBM Plex Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    bg._element.spPr.append(bg._element.spPr.makeelement(qn('a:effectLst'), {}))
    return slide


def add_text(slide, left, top, width, height, text, size=18, color=TEXT, bold=False,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
             letter_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return box


def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def eyebrow(slide, text, top=Inches(0.7)):
    add_rect(slide, Inches(0.9), top + Inches(0.12), Inches(0.3), Pt(1.4), fill=ACCENT)
    add_text(slide, Inches(1.3), top, Inches(8), Inches(0.35), text.upper(),
              size=13, color=ACCENT, bold=True, font=FONT_MONO)


def headline(slide, text, top=Inches(1.15), size=38, width=Inches(11)):
    add_text(slide, Inches(0.9), top, width, Inches(1.4), text, size=size, color=TEXT,
              bold=True, font=FONT_BODY, line_spacing=1.0)


def slide_number(slide, n):
    add_text(slide, SW - Inches(1.6), SH - Inches(0.9), Inches(1.3), Inches(0.6),
              str(n).zfill(2), size=44, color=SURFACE2, bold=True, font=FONT_MONO,
              align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, tag, title, desc, tag_color=ACCENT):
    add_rect(slide, left, top, width, height, fill=SURFACE, line_color=LINE)
    pad = Inches(0.2)
    add_text(slide, left + pad, top + Inches(0.16), width - 2 * pad, Inches(0.3), tag,
              size=10.5, color=tag_color, bold=True, font=FONT_MONO)
    add_text(slide, left + pad, top + Inches(0.52), width - 2 * pad, Inches(0.4), title,
              size=14.5, color=TEXT, bold=True, font=FONT_BODY)
    add_text(slide, left + pad, top + Inches(0.92), width - 2 * pad, height - Inches(1.1), desc,
              size=11.5, color=TEXT_DIM, font=FONT_BODY, line_spacing=1.2)


def grid_cards(slide, items, cols, top=Inches(1.9), left=Inches(0.9),
               total_width=Inches(11.5), row_h=Inches(1.65), gap=Inches(0.12)):
    col_w = (total_width - gap * (cols - 1)) / cols
    for i, (tag, title, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = left + c * (col_w + gap)
        y = top + r * (row_h + gap)
        card(slide, x, y, col_w, row_h, tag, title, desc)


def bullet_list(slide, items, top=Inches(1.9), left=Inches(0.9), width=Inches(10.5),
                 size=15, gap=Inches(0.62)):
    for i, (strong, rest) in enumerate(items):
        y = top + i * gap
        add_rect(slide, left, y + Inches(0.06), Inches(0.09), Inches(0.09), fill=ACCENT)
        add_text(slide, left + Inches(0.25), y, Inches(0.15), Inches(0.4), "",
                  size=size, color=TEXT_DIM)
        box = slide.shapes.add_textbox(left + Inches(0.25), y - Inches(0.08), width, Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = strong + "  "
        r1.font.bold = True
        r1.font.size = Pt(size)
        r1.font.color.rgb = TEXT
        r1.font.name = FONT_BODY
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.color.rgb = TEXT_DIM
        r2.font.name = FONT_BODY
        add_rect(slide, left, y + Inches(0.5), width + Inches(0.25), Pt(1), fill=LINE)


def tag_pill(slide, x, y, text, on=False):
    w = Inches(0.18) + Inches(0.095) * len(text)
    fill = RGBColor(0x2A, 0x1E, 0x10) if on else SURFACE
    line = ACCENT_SOFT if on else LINE_STRONG
    color = ACCENT if on else TEXT_DIM
    add_rect(slide, x, y, w, Inches(0.34), fill=fill, line_color=line)
    add_text(slide, x, y + Inches(0.055), w, Inches(0.3), text, size=10.5, color=color,
              font=FONT_MONO, align=PP_ALIGN.CENTER)
    return w


def tag_row(slide, tags, top, left=Inches(0.9)):
    x = left
    for text, on in tags:
        w = tag_pill(slide, x, top, text, on)
        x += w + Inches(0.14)


def filetree(slide, lines, top=Inches(1.9), height=Inches(5.0)):
    add_rect(slide, Inches(0.9), top, Inches(11.5), height, fill=SURFACE, line_color=LINE)
    box = slide.shapes.add_textbox(Inches(1.25), top + Inches(0.28), Inches(10.9),
                                    height - Inches(0.56))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (indent, runs) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.4
        if indent:
            r0 = p.add_run()
            r0.text = "    " * indent
            r0.font.size = Pt(13.5)
            r0.font.name = FONT_MONO
        for color, text in runs:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(13.5)
            run.font.name = FONT_MONO
            run.font.color.rgb = color


# ---------------------------------------------------------------- Slide 1
s = add_slide()
add_text(s, 0, Inches(0.9), SW, Inches(0.4), "PROJET WEB TEMPS REEL", size=13, color=ACCENT,
          bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(1.5), SW, Inches(1.6), "PONG MULTIJOUEUR", size=64, color=TEXT, bold=True,
          font=FONT_BODY, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.5), Inches(2.9), Inches(8.33), Inches(1.0),
          "Un Pong jouable dans le navigateur, en temps reel, ou deux joueurs s'affrontent "
          "via WebSocket — toute la physique arbitree cote serveur.",
          size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER, line_spacing=1.3)
tag_row(s, [("Node.js · ws", False), ("Canvas HTML5", False), ("WebSocket", False)],
        Inches(4.1), left=Inches(4.9))

# ---------------------------------------------------------------- Slide 2
s = add_slide()
eyebrow(s, "Contexte")
headline(s, "Un jeu, une contrainte : le temps reel")
add_text(s, Inches(0.9), Inches(2.3), Inches(10.8), Inches(1.0),
          "L'objectif est de developper un Pong jouable dans un navigateur, permettant a deux "
          "joueurs de se connecter simultanement et de s'affronter via une connexion WebSocket "
          "— sans latence perceptible, sans triche possible.",
          size=17, color=TEXT_DIM, line_spacing=1.4)
add_text(s, Inches(0.9), Inches(3.4), Inches(10.8), Inches(1.0),
          "Le serveur Node.js fait autorite sur toute la logique de jeu ; le client se limite a "
          "afficher l'etat recu et a transmettre les touches pressees.",
          size=17, color=TEXT_DIM, line_spacing=1.4)
slide_number(s, 2)

# ---------------------------------------------------------------- Slide 3
s = add_slide()
eyebrow(s, "Cahier des charges")
headline(s, "Objectifs fonctionnels")
grid_cards(s, [
    ("01", "Deux joueurs", "Rejoindre une meme partie, a minima."),
    ("02", "Plateau affiche", "Terrain, balle et deux raquettes rendus au canvas."),
    ("03", "Synchronisation", "Positions des raquettes et de la balle en temps reel."),
    ("04", "Score", "Suivi du score de chaque joueur."),
    ("05", "Fin de partie", "Detectee par deconnexion ou score maximum atteint."),
    ("06", "Salons de jeu", "Plusieurs parties isolees, en parallele."),
], cols=3)
slide_number(s, 3)

# ---------------------------------------------------------------- Slide 4
s = add_slide()
eyebrow(s, "Au-dela du socle")
headline(s, "Fonctionnalites optionnelles")
grid_cards(s, [
    ("IMPLEMENTE ✓", "Mode spectateur",
     "Rejoindre une partie en cours pour la regarder, sans y participer — aucun input "
     "envoye, reception seule de l'etat de jeu."),
    ("IMPLEMENTE ✓", "Classement des joueurs",
     "Historique victoires / defaites tenu en memoire cote serveur, consultable a tout "
     "moment depuis le lobby."),
], cols=2, row_h=Inches(2.2))
slide_number(s, 4)

# ---------------------------------------------------------------- Slide 5
s = add_slide()
eyebrow(s, "Exigences")
headline(s, "Contraintes techniques")
bullet_list(s, [
    ("Latence", "— le jeu doit rester jouable malgre le reseau"),
    ("Autorite serveur", "— toute la logique calculee cote serveur, pour eviter la triche"),
    ("Scalabilite", "— gerer plusieurs parties en parallele"),
    ("Robustesse", "— deconnexions/reconnexions gerees sans planter la partie"),
    ("Compatibilite", "— navigateurs modernes (Chrome, Firefox…)"),
])
slide_number(s, 5)

# ---------------------------------------------------------------- Slide 6
s = add_slide()
eyebrow(s, "Architecture")
headline(s, "Client leger, serveur autoritaire")
top = Inches(2.3)
bw = Inches(4.6)
card(s, Inches(0.9), top, bw, Inches(3.6), "CLIENT", "Canvas HTML5",
     "• Affiche l'etat recu du serveur\n• Capture les touches (haut / bas)\n"
     "• Envoie les inputs, jamais la physique", tag_color=INFO)
card(s, Inches(8.0), top, bw, Inches(3.6), "SERVEUR", "Node.js + ws",
     "• Calcule la physique (balle, collisions)\n• Arbitre le score et la fin de partie\n"
     "• Diffuse l'etat a 60 Hz a tous les clients", tag_color=ACCENT)
add_text(s, Inches(5.6), top + Inches(1.4), Inches(2.3), Inches(0.4), "WEBSOCKET",
          size=11, color=TEXT_DIM, font=FONT_MONO, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.6), top + Inches(1.85), Inches(2.3), Pt(1.4), fill=LINE_STRONG)
add_text(s, Inches(5.6), top + Inches(1.95), Inches(2.3), Inches(0.4), "JSON, ~20/s",
          size=11, color=TEXT_DIM, font=FONT_MONO, align=PP_ALIGN.CENTER)
slide_number(s, 6)

# ---------------------------------------------------------------- Slide 7
s = add_slide()
eyebrow(s, "Cycle de vie")
headline(s, "Du clic « Jouer » au point marque")
steps = [
    ("1", "Connexion", "Le joueur se connecte, le serveur l'affecte a un salon disponible ou en cree un."),
    ("2", "Salle complete", "Des que deux joueurs sont presents dans le salon, la partie demarre."),
    ("3", "Boucle de jeu", "Le client envoie ses inputs ; le serveur met a jour l'etat et le diffuse a tous."),
    ("4", "Fin de partie", "Score maximum atteint ou deconnexion d'un joueur : la partie se termine."),
]
top = Inches(2.1)
for i, (n, title, desc) in enumerate(steps):
    y = top + i * Inches(1.05)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y, Inches(0.4), Inches(0.4))
    circ.fill.solid(); circ.fill.fore_color.rgb = BG
    circ.line.color.rgb = ACCENT_SOFT; circ.line.width = Pt(1)
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run(); run.text = n
    run.font.size = Pt(13); run.font.color.rgb = ACCENT; run.font.name = FONT_MONO
    add_text(s, Inches(1.55), y - Inches(0.02), Inches(9.8), Inches(0.35), title,
              size=17, color=TEXT, bold=True)
    add_text(s, Inches(1.55), y + Inches(0.34), Inches(9.8), Inches(0.55), desc,
              size=12.5, color=TEXT_DIM, line_spacing=1.25)
slide_number(s, 7)

# ---------------------------------------------------------------- Slide 8
s = add_slide()
eyebrow(s, "Methode")
headline(s, "Plan de developpement")
grid_cards(s, [
    ("01", "Serveur WebSocket basique", "Connexion / deconnexion."),
    ("02", "Rooms a 2 joueurs", "Matchmaking automatique."),
    ("03", "Rendu canvas statique", "Terrain, raquettes, balle."),
    ("04", "Inputs clavier", "Envoi au serveur."),
    ("05", "Game loop serveur", "Physique, collisions, score."),
    ("06", "Diffusion + rendu dynamique", "Synchronisation temps reel."),
    ("07", "Fin de partie", "Detection et redemarrage."),
    ("08", "Ameliorations", "Spectateurs, classement."),
], cols=4, row_h=Inches(1.9))
slide_number(s, 8)

# ---------------------------------------------------------------- Slide 9
s = add_slide()
eyebrow(s, "Code · Serveur")
headline(s, "server/")
filetree(s, [
    (0, [(TEXT, "index.js"), (TEXT_FAINT, "  — HTTP statique + WebSocket, routage des messages")]),
    (1, [(TEXT, "lib/")]),
    (2, [(ACCENT, "Paddle.js, Ball.js"), (TEXT_FAINT, "  — physique elementaire")]),
    (2, [(ACCENT, "GameState.js"), (TEXT_FAINT, "  — collisions, score, snapshot")]),
    (2, [(ACCENT, "Room.js, RoomManager.js"), (TEXT_FAINT, "  — salons, matchmaking")]),
    (2, [(ACCENT, "Connection.js"), (TEXT_FAINT, "  — wrapper autour du WebSocket brut")]),
    (2, [(ACCENT, "LeaderboardStore.js"), (TEXT_FAINT, "  — victoires / defaites")]),
    (2, [(ACCENT, "constants.js, messageTypes.js")]),
], height=Inches(4.3))
slide_number(s, 9)

# ---------------------------------------------------------------- Slide 10
s = add_slide()
eyebrow(s, "Code · Client")
headline(s, "client/")
filetree(s, [
    (0, [(TEXT, "index.html, css/style.css"), (TEXT_FAINT, "  — lobby, jeu, fin de partie, classement")]),
    (1, [(TEXT, "js/")]),
    (2, [(ACCENT, "network.js"), (TEXT_FAINT, "  — connexion WebSocket, envoi/reception typee")]),
    (2, [(ACCENT, "state.js"), (TEXT_FAINT, "  — etat partage cote client")]),
    (2, [(ACCENT, "renderer.js"), (TEXT_FAINT, "  — rendu canvas statique puis dynamique")]),
    (2, [(ACCENT, "input.js"), (TEXT_FAINT, "  — clavier + tactile mobile")]),
    (2, [(ACCENT, "ui.js"), (TEXT_FAINT, "  — ecrans, score, classement")]),
    (2, [(ACCENT, "main.js"), (TEXT_FAINT, "  — cablage des evenements")]),
], height=Inches(4.3))
slide_number(s, 10)

# ---------------------------------------------------------------- Slide 11
s = add_slide()
eyebrow(s, "Coeur du serveur")
headline(s, "La boucle de jeu, 60 fois par seconde")
bullet_list(s, [
    ("•", "Mise a jour des raquettes selon l'input recu"),
    ("•", "Deplacement de la balle et collision avec les murs"),
    ("•", "Rebond sur raquette : angle selon le point d'impact"),
    ("•", "Vitesse incrementee a chaque echange"),
    ("•", "Diffusion de l'etat (JSON) a tous les clients du salon"),
], left=Inches(0.9), width=Inches(6.0), size=13.5, gap=Inches(0.52))
card(s, Inches(7.4), Inches(2.3), Inches(5.0), Inches(2.6), "ANTI-TRICHE",
     "Aucune physique cote client",
     "Le client n'envoie que { up, down }. Position de la balle, collisions et score "
     "n'existent que cote serveur — impossible a falsifier depuis le navigateur.")
slide_number(s, 11)

# ---------------------------------------------------------------- Slide 12
s = add_slide()
eyebrow(s, "Multi-parties")
headline(s, "Salons & matchmaking")
add_text(s, Inches(0.9), Inches(2.3), Inches(11), Inches(1.4),
          "RoomManager.findOrCreateRoom() place chaque nouveau joueur dans un salon existant "
          "non complet, ou en cree un nouveau — chaque salon possede son propre GameState, "
          "totalement independant des autres.",
          size=17, color=TEXT_DIM, line_spacing=1.4)
tag_row(s, [("Parties en parallele", True), ("Etat isole par salon", True),
            ("Nettoyage a la deconnexion", True)], Inches(3.7))
slide_number(s, 12)

# ---------------------------------------------------------------- Slide 13
s = add_slide()
eyebrow(s, "Fonctionnalites optionnelles")
headline(s, "Spectateurs & classement")
grid_cards(s, [
    ("join_spectate", "Regarder une partie",
     "Rejoint un salon deja en cours de jeu, recoit les memes diffusions d'etat que les "
     "joueurs, sans jamais emettre d'input."),
    ("LeaderboardStore", "Victoires / defaites",
     "Chaque fin de partie incremente le score du vainqueur et du perdant ; le classement "
     "est trie par victoires."),
], cols=2, row_h=Inches(2.2))
slide_number(s, 13)

# ---------------------------------------------------------------- Slide 14
s = add_slide()
eyebrow(s, "Retour d'experience")
headline(s, "Ce qui a resiste : le reseau, pas le code", size=32)
add_text(s, Inches(0.9), Inches(2.05), Inches(11), Inches(0.8),
          "Le jeu tournait des les premiers tests, mais le rendre jouable entre un PC et un "
          "telephone sur le meme Wi-Fi a exige un vrai diagnostic reseau.",
          size=14.5, color=TEXT_DIM, line_spacing=1.3)
grid_cards(s, [
    ("OBSTACLE", "NAT WSL2",
     "Le serveur, lance dans WSL, restait invisible depuis le reseau local malgre le port forwarding."),
    ("OBSTACLE", "Pare-feu tiers",
     "Un antivirus avec pare-feu propre bloquait les connexions entrantes non reconnues."),
    ("OBSTACLE", "Pas de clavier",
     "Les controles clavier ne fonctionnaient evidemment pas sur mobile."),
], cols=3, top=Inches(3.0), row_h=Inches(1.5))
tag_row(s, [("→ Execution native Windows", True), ("→ Regles de pare-feu dediees", True),
            ("→ Controles tactiles ajoutes", True)], Inches(4.85))
slide_number(s, 14)

# ---------------------------------------------------------------- Slide 15
s = add_slide()
eyebrow(s, "Bilan")
headline(s, "Le coeur du jeu, et au-dela")
add_text(s, Inches(0.9), Inches(2.15), Inches(11), Inches(1.0),
          "Socle fonctionnel complet — connexion, salons, synchronisation, score, fin de "
          "partie — auquel s'ajoutent les deux fonctionnalites optionnelles : spectateurs "
          "et classement.",
          size=16, color=TEXT_DIM, line_spacing=1.35)
add_text(s, Inches(0.9), Inches(3.3), Inches(4), Inches(0.35), "PERSPECTIVES",
          size=12.5, color=ACCENT, bold=True, font=FONT_MONO)
tag_row(s, [("Reconnexion apres coupure", False), ("Deploiement public (HTTPS)", False)],
        Inches(3.75))
tag_row(s, [("Tournois a elimination", False), ("Historique de parties detaille", False)],
        Inches(4.25))
slide_number(s, 15)

prs.save("/home/bastianbondoux/Pong-project/pong-deck.pptx")
print("OK")
